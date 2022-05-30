import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import openpyxl
from images_loader import SheetImageLoader
import os
import io
from PIL import Image
from svg2png_svglib import convertSVG2PNG
import sys


def get_images(cell, image_loader):
    images = image_loader.get(cell)
    return images


def get_column_letter(column_index):
    """Return the column letter for a column index.
    For example, column_index 3 would return 'C'.
    """
    column_letter = ""
    while column_index > 0:
        column_letter = chr(ord('A') + (column_index-1) % 26) + column_letter
        column_index = int((column_index-1) / 26)
    return column_letter


def getExcelData(workbook, sheetName):
    pxl_doc = openpyxl.load_workbook(workbook)
    sheet = pxl_doc[sheetName]
    image_loader = SheetImageLoader(sheet)
    excelData = {}
    # iterate over columns in the first row
    for i in range(1, sheet.max_column+1):
        # get column letter
        col = get_column_letter(i)
        images = get_images(col+'2', image_loader)
        if len(images) == 0:
            excelData[sheet.cell(row=1, column=i).value] = sheet.cell(row=2, column=i).value
            # print(sheet.cell(row=1, column=i).value, ", ", sheet.cell(row=2, column=i).value)
        else:
            excelData[sheet.cell(row=1, column=i).value] = images
            # print(sheet.cell(row=1, column=i).value, ", n° of images: ", len(images))
    return excelData


data = getExcelData("./NPI_TEMPLATE_FILL.xlsx", 'TEMPLATE_1_FILL')
print("Data Keys =========")
print(data.keys())
print("Data Keys ========= END")
# exit()




# dataKeys = ['Template', 'Template Colors', 'Title', 'Title Description', 'Supplier', 'Background Image', 'Why it matters?', 'Key Benefits 1', 'Key Benefits 2', 'Alt Table Title', 'Alt table Text', 'Product Image 1', 'Product Image 2', 'Product Image 3', 'OPN 1', 'Description 1', 'OPN 1 Link', 'OPN 2', 'Description 2', 'OPN 2 Link', 'Application 1 Title', 'Application 1 Background', 'Application 1 Description', 'Application 2 Title', 'Application 2 Background', 'Application 2 Description', 'Application 3 Title', 'Application 3 Background', 'Application 3 Description', 'Application 4 Title', 'Application 4 Background', 'Application 4 Description', 'TDK MEMS?']
# dataKeys = ['#', 'Template', 'Template Colors', 'Title', 'Title Description', 'Supplier', 'Background Image', 'Why it matters?', 'Key Benefits 1', 'Key Benefits 2', 'Alt Table Title', 'Alt table Text', 'Product Image 1', 'Product Image 2', 'Product Image 3', 'OPN 1 Name', 'OPN 1 Description', 'OPN 1 Link', 'OPN 2 Name', 'OPN 2 Description', 'OPN 2 Link', 'Application 1 Title', 'Application 1 Background', 'Application 1 Description', 'Application 2 Title', 'Application 2 Background', 'Application 2 Description', 'Application 3 Title', 'Application 3 Background', 'Application 3 Description', 'Application 4 Title', 'Application 4 Background', 'Application 4 Description', 'TDK MEMS?', 'Feature 1 name', 'Feature 2 name', 'Feature 1 rating', 'Feature 2 rating']

# get data keys from first row of excel sheet
dataKeys = list(data.keys())
# in general
applicationIcons = []
for key in dataKeys:
    if key.startswith("Application"):
        if key.endswith("Icon"):
            applicationIcons.append(data[key])

applicationBackgrounds = []
for key in dataKeys:
    if key.startswith("Application"):
        if key.endswith("Background"):
            applicationBackgrounds.append(data[key])

applicationTitles = []
for key in dataKeys:
    if key.startswith("Application"):
        if key.endswith("Title"):
            applicationTitles.append(data[key])

applicationDescriptions = []
for key in dataKeys:
    if key.startswith("Application"):
        if key.endswith("Description"):
            applicationDescriptions.append(data[key])

OPNNames = []
for key in dataKeys:
    if key.startswith("OPN"):
        if key.endswith("Name"):
            OPNNames.append(data[key])

OPNDescriptions = []
for key in dataKeys:
    if key.startswith("OPN"):
        if key.endswith("Description"):
            OPNDescriptions.append(data[key])

OPNLinks = []
for key in dataKeys:
    if key.startswith("OPN"):
        if key.endswith("Link"):
            OPNLinks.append(data[key])

featureNames = []
for key in dataKeys:
    if key.startswith("Feature"):
        if key.endswith("name"):
            featureNames.append(data[key])

featureRatings = []
for key in dataKeys:
    if key.startswith("Feature"):
        if key.endswith("rating"):
            featureRatings.append(data[key])

tables1 = {
    "application": {
        "icons": applicationIcons,
        "titles": applicationTitles,
        "descriptions": applicationDescriptions
    },
    "OPN": {
        #"names" = OPNNames+OPNLinks per item
        "names": OPNNames,  # [f"{OPNNames[i]}{{{OPNLinks[i]}}}" for i in range(len(OPNNames))],
        "descriptions": OPNDescriptions,
        "links": [f"LINK{i+1}" for i in range(len(OPNNames))]
    },
    "feature": {
        "names": featureNames,
        "ratings": featureRatings
    }
}
hyperlinks1 = {
    "OPN": {
        "names": OPNLinks,
        "links": OPNLinks
    }
}
backgrounds1 = {
    "application": {
        "titles": applicationBackgrounds,
    }
}


def fillPPTXTable(tables, tableName, shape, shapes):
    tableDict = tables[tableName]
    # set dict keys as column names
    tkeys = list(tableDict.keys())
    # for i in range(len(tkeys)):
    #     print(i)
    #     table.cell(0, i-1).text = tkeys[i-1]
    # set dict values as column values
    numRows = len(list(tableDict.values())[0])
    numCols = len(tkeys)

    if shape.has_table:
        table = shape.table
    else:
        if shape.has_text_frame:
            shape.text_frame.clear()
        tableSize = shape.width, shape.height

        # for column in shape.table.columns:
        #     column.width = 0
        # sp = shape._sp
        # sp.getparent().remove(sp)
        table = shapes.add_table(numRows+1, numCols, shape.left, shape.top, *tableSize).table
    for i in range(len(tkeys)):
        # table.cell(j+1, i-1).text = str(tableDict[tkeys[i-1]][j])
        text_frame = table.cell(0, i-1).text_frame
        text_frame.clear()  # not necessary for newly-created shape

        p = text_frame.paragraphs[0]
        # p.text = str(tableDict[tkeys[i-1]][j])
        run = p.add_run()
        run.text = tkeys[i-1].capitalize()
        # if tableName in hyperlinks1.keys():
        #     if tkeys[i-1] in hyperlinks1[tableName].keys():
        #         run.hyperlink.address = hyperlinks1[tableName][tkeys[i-1]][j]
        font = run.font
        font.name = 'Arrow Display'
        font.size = Pt(10)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        # set color to black
        font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        # align text to center
        p.alignment = PP_ALIGN.CENTER
        for j in range(numRows):
            # table.cell(j+1, i-1).text = str(tableDict[tkeys[i-1]][j])
            text_frame = table.cell(j+1, i-1).text_frame
            text_frame.clear()  # not necessary for newly-created shape

            p = text_frame.paragraphs[0]
            # p.text = str(tableDict[tkeys[i-1]][j])
            run = p.add_run()
            run.text = str(tableDict[tkeys[i-1]][j])
            if tableName in hyperlinks1.keys():
                if tkeys[i-1] in hyperlinks1[tableName].keys():
                    run.hyperlink.address = hyperlinks1[tableName][tkeys[i-1]][j]
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(7)
            font.bold = True
            font.italic = None  # cause value to be inherited from theme
            # set color to black
            font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            # align text to center
            p.alignment = PP_ALIGN.CENTER

# dataKeys = ['#', 'Template', 'Template Colors', 'Title', 'Title Description', 'Supplier', 'Background Image', 'Why it matters?', 'Key Benefits 1', 'Key Benefits 2', 'Alt Table Title', 'Alt table Text', 'Product Image 1', 'Product Image 2', 'Product Image 3', 'OPN 1 Name', 'OPN 1 Description', 'OPN 1 Link', 'OPN 2 Name', 'OPN 2 Description', 'OPN 2 Link', 'Application 1 Title', 'Application 1 Background', 'Application 1 Description', 'Application 2 Title', 'Application 2 Background', 'Application 2 Description', 'Application 3 Title', 'Application 3 Background', 'Application 3 Description', 'Application 4 Title', 'Application 4 Background', 'Application 4 Description', 'TDK MEMS?', 'Feature 1 name', 'Feature 2 name', 'Feature 1 rating', 'Feature 2 rating']




def addRun(paragraph, text, font_size=20, bold=False, italic=False, underline=False, color=None):
    run = paragraph.add_run()
    run.text = str(text)
    font = run.font
    font.name = 'Arrow Display'
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.underline = underline
    if color is not None:
        font.color.rgb = color


def writeText(text, text_frame, font_size=15, bold=False, italic=False, underline=False, color=None, dotted=False, replace=False, append=False):
    if not replace and not append:
        text_frame.clear()  # not necessary for newly-created shape
    if not dotted:
        if not append:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        if replace:
            p.text = text
        else:
            addRun(p, text, font_size, bold, italic, underline, color)
    else:
        if not append:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        if isinstance(text, list):
            paragraph_strs = text
        elif isinstance(text, str):
            paragraph_strs = text.split("\n")
        addRun(p, paragraph_strs[0], font_size, bold, italic, underline, color)
        for para_str in paragraph_strs[1:]:
            p = text_frame.add_paragraph()
            addRun(p, para_str, font_size=font_size, bold=bold, italic=italic, underline=underline, color=color)


def putImage(imgName, shape, shapes):
    """
    Adds an image that fills the entire shape.
    """
    shapeDimensions = shape.width, shape.height
    # convert emu to pixels
    shapeWidth = int(shapeDimensions[0] / 9525)
    shapeHeight = int(shapeDimensions[1] / 9525)
    assert shapeWidth < 3000
    assert shapeHeight < 2000
    shapeLocation = shape.left, shape.top

    # delete shape
    sp = shape._sp
    sp.getparent().remove(sp)
    imgPath = os.path.join(os.path.dirname(__file__), "images", imgName)

    # check if svg
    if imgName.endswith(".svg"):
        # check if png version exists
        png_path = imgPath[:-4] + ".png"
        if os.path.exists(png_path):
            imgPath = png_path
        else:
            convertSVG2PNG(imgPath, png_path)
            imgPath = png_path
    # create shape in the same location as the original shape
    print("Invoked from putImage {}".format(shapeLocation))
    new_shape = shapes.add_picture(imgPath, shapeLocation[0], shapeLocation[1])

    # fit image to shape
    new_shape.width = shapeDimensions[0]
    new_shape.height = shapeDimensions[1]


def putImages(imgNames, shape, shapes):
    """
    Adds images equidistant from each other along the shape height
    """
    shapeDimensions = shape.width, shape.height
    # convert emu to pixels
    shapeWidth = int(shapeDimensions[0] / 9525)
    shapeHeight = int(shapeDimensions[1] / 9525)
    assert shapeWidth < 3000
    assert shapeHeight < 2000
    shapeLocation = shape.left, shape.top
    # delete shape
    sp = shape._sp
    sp.getparent().remove(sp)
    imgPaths = [os.path.join(os.path.dirname(__file__), "images", imgName) for imgName in imgNames if imgName != ""]
    imagesHeight = int(shapeHeight / len(imgPaths))
    for i, imgPath in enumerate(imgPaths):
        # check if svg
        if imgPath.endswith(".svg"):
            # check if png version exists
            png_path = imgPath[:-4] + ".png"
            if os.path.exists(png_path):
                imgPath = png_path
            else:
                convertSVG2PNG(imgPath, png_path)
                imgPath = png_path
        # create shape in the same location as the original shape
        print("Invoked from putImages {}".format(shapeLocation))
        new_shape = shapes.add_picture(imgPath, shapeLocation[0], shapeLocation[1] + imagesHeight * i * 9525)

        # fit image to shape
        new_shape.width = shapeDimensions[0]
        new_shape.height = imagesHeight*9525


def writeEquidistantTexts(texts, shape, shapes, backgrounds=None):
    """
    Adds texts equidistant from each other along the shape height
    """
    # TODO add backgrounds
    shapeDimensions = shape.width, shape.height
    # convert emu to pixels
    shapeWidth = int(shapeDimensions[0] / 9525)
    shapeHeight = int(shapeDimensions[1] / 9525)
    assert shapeWidth < 3000
    assert shapeHeight < 2000
    shapeLocation = shape.left, shape.top
    if backgrounds is not None:
        putImages(backgrounds, shape, shapes)
        print(backgrounds)
    else:
        # delete shape
        sp = shape._sp
        sp.getparent().remove(sp)
    textsHeight = int(shapeHeight / len(texts))
    for i, text in enumerate(texts):
        # create shape in the same location as the original shape
        new_shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, shapeLocation[0], shapeLocation[1] + textsHeight * i * 9525, shapeWidth*9525, textsHeight*9525)
        # add transparent fill
        new_shape.fill.background()
        # write text
        writeText(text, new_shape.text_frame)


shapes_1 = {
    "0": {
        "1": "Title",
        "2": "Supplier",
        "4": "img:Background Image",
    },
    "1": {
        "0": "Why it matters?",
        "1": ["Key Benefits 1", "Key Benefits 2"],
        # "2": "Subject",
        "5": "img:Supplier Image",
        "4": "Product Image 1",
    },
    "2": {
        "1": "imgs:application:icons",
        "5": "application:titles",
        "8": "img:Supplier Image",
    },
    "3": {}
}

table_shape_1 = {
    "0": {},
    "1": {},
    "2": {
        "7": "feature",
        "0": "OPN",
        # "8": "application"
    },
    "3": {}
}

specialFormat_1 = {
    "Key Benefits 1": {"bold": True, "dotted": True},
    "Key Benefits 2": {"bold": True, "dotted": True, "append": True},
    # set title color white
    "Title": {"bold": True, "font_size": 30, "color": RGBColor(0xFF, 0xFF, 0xFF)},
    "Supplier": {"bold": True, "font_size": 20, "color": RGBColor(0xFF, 0xFF, 0xFF)},
}

"""
Supplier --> Title
Topic- Part Number --> Title Description
Overview -> Why it matters?
Features -> Key Benefits 1
Benefits -> Key Benefits 2
Specs Table -> Idéntica a la tabla construida para template 1 (adjuntando imagen)
Market & Applications -> Application 1 Title, Application 2 Title, Application 3 Title y Application 4 Title (en lugar de backgrounds habría íconos que estarían a la derecha)
"""

table_shape_2 = {
    "0": {},
    "1": {},
    "2": {
        "7": "feature",
        "8": "OPN"
    },
    "3": {}
}

shapes_2 = {
    "0":
        {
            "2": f"{data['Supplier']}\n{data['Title Description']}",
        },
    "1":
        {
            "1": data["Title Description"],
            "2": "Key Benefits 1",
            "5": "Key Benefits 2",
            "6": "Product Image 1",
            "7": "img:Supplier Image",
        },
    "2":
        {
            "5": "application:titles",
            "9": "img:Supplier Image",
            "13": "img:Application 1 Icon",
            "10": "img:Application 2 Icon",
            "11": "img:Application 3 Icon",
            "12": "img:Application 4 Icon",
        },
    "3":
        {}
    }

table_shape_3 = {
    '0': {},
    '1': {},
    '2': {
        '1': 'feature',
        '6': 'OPN',

    }
}

shapes_3 = {
    '0': {
        '3': "img:Background Image",
        '2': "Supplier",
        '4': "Title",
    },
    '1': {
        '4': 'Overview',
        '6': 'Key Benefits 1',
        '7': 'Key Benefits 2',
        '8': "img:Supplier Image",
        '5': "Product Image 1",
    },
    '2': {
        '8': 'img:Application 1 Icon',
        '9': 'img:Application 2 Icon',
        '10': 'img:Application 3 Icon',
        '11': 'img:Application 4 Icon',
        '5': 'application:titles',
        '7': 'img:Supplier Image',
    }
}


# template = int(data["Template"])
##template = int(sys.argv[1])
template = 2
prs = Presentation(f'Template {template}.pptx')

if template == 1:
    shapes = shapes_1
    tables_shapes = table_shape_1
    tables = tables1
    backgrounds = backgrounds1
elif template == 2:
    shapes = shapes_2
    tables_shapes = table_shape_2
    tables = tables1
    backgrounds = {}
elif template == 3:
    shapes = shapes_3
    tables_shapes = table_shape_3
    tables = tables1
    backgrounds = {}

for i, slide in enumerate(prs.slides):
    # print(i, slide)
    slideDict = shapes[str(i)]
    slideTableDict = tables_shapes[str(i)]
    for j, shape in enumerate(slide.shapes):
        # if shape.has_table:
        if str(j) in slideTableDict.keys():
            # table = shape.table
            tableName = slideTableDict[str(j)]
            # print(i, j)
            fillPPTXTable(tables, tableName, shape, slide.shapes)
        elif shape.has_text_frame:
            if template == 3 and j == 4 and i == 2:
                # clear text
                shape.text_frame.clear()
                paragraph = shape.text_frame.paragraphs[0]
                # add a run
                run = paragraph.add_run()
                run.text = "Additional Information:\n"
                run = paragraph.add_run()
                # set the text of the run containing an hiperlink
                run.text = 'Dataset: Link'
                # set the hyperlink
                run.hyperlink.address = data['OPN 1 Link']
            if str(j) in slideDict.keys():
                # writeText(str(slideDict[str(j)]), shape.text_frame)
                inp = slideDict[str(j)]
                if inp in dataKeys:
                    # if data[inp] is string or number
                    if isinstance(data[inp], str) or isinstance(data[inp], int) or isinstance(data[inp], float):
                        if inp in specialFormat_1.keys():
                            writeText(str(data[inp]), shape.text_frame, **specialFormat_1[inp])
                        else:
                            writeText(str(data[inp]), shape.text_frame, replace=True)
                    else:
                        image = data[inp][0]
                        with io.BytesIO() as output:
                            # print(image)
                            image.save(output, format="GIF")
                            left = shape.left
                            top = shape.top
                            pic = slide.shapes.add_picture(output, left, top, height=shape.height, width=shape.width)
                else:
                    if isinstance(inp, list):
                        shape.text_frame.clear()
                        for m in inp:
                            if m in specialFormat_1.keys():
                                writeText(str(data[m]), shape.text_frame, **specialFormat_1[m])
                            else:
                                writeText(data[m], shape.text_frame)
                        # writeText("Key Features:", shape.text_frame, font_size=12, bold=True, replace=True)
                        # writeText(data[inp[0]], shape.text_frame, append=True, bold=True, font_size=15, replace=False, dotted=False)
                        # writeText(data[inp[1]], shape.text_frame, append=True, font_size=15, replace=False, dotted=False)

                    # chech if ":" in inp string
                    elif inp.startswith("img:"):
                        # print(inp)
                        rinp = inp[4:]
                        if rinp in dataKeys:
                            # print(rinp, data[rinp])
                            putImage(data[rinp], shape, slide.shapes)
                            # if rinp in specialFormat_1.keys():
                            #     writeText("img:"+str(data[rinp]), shape.text_frame, **specialFormat_1[rinp])
                            # else:
                            #     writeText("img:"+str(data[rinp]), shape.text_frame, replace=True)
                        elif ":" in rinp:
                            table, column = rinp.split(":")
                            if table in tables.keys():
                                if column in tables[table].keys():
                                    for img in tables[table][column]:
                                        putImage(img, shape, slide.shapes)
                                    # writeText(tables[table][column], shape.text_frame, dotted=True)
                    elif inp.startswith("imgs:"):
                        rinp = inp[5:]
                        if ":" in rinp:
                            table, column = rinp.split(":")
                            if table in tables.keys():
                                if column in tables[table].keys():
                                    putImages(tables[table][column], shape, slide.shapes)
                                    # for img in tables[table][column]:
                                    #     putImage(img, shape, slide.shapes)
                                    # writeText(tables[table][column], shape.text_frame, dotted=True)
                    # elif inp has ":"
                    elif ":" in inp:
                        table, column = inp.split(":")
                        if table in tables.keys():
                            # if column in tables[table].keys():
                            bgs = None
                            if table in backgrounds.keys():
                                if column in backgrounds[table].keys():
                                    bgs = backgrounds[table][column]
                            writeEquidistantTexts(tables[table][column], shape, slide.shapes, bgs)
                            # writeText(tables[table][column], shape.text_frame, dotted=True)
                    else:
                        writeText(str(inp), shape.text_frame)
            # else:
            #     print(i, j)
            #     if shape.has_text_frame:
            #         writeText(str(j), shape.text_frame)

prs.save(f'test {template}.pptx')
print(f"saved at test {template}.pptx")
