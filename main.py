import collections
import collections.abc
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
import pptx.enum.shapes as pptx_types
from PIL import Image

from svg2png_svglib import convertSVG2PNG

import openpyxl

MAX_ROWS = 6;
global_AUX = None
Separation_Char = None;
EmptySpace_Char = None;

os.system("cls")

def printWARN(text):
    print("\033[33m{}\033[0m".format(text));

def get_column_letter(column_index):
    """Return the column letter for a column index.
    For example, column_index 3 would return 'C'.
    """
    column_letter = ""
    while column_index > 0:
        column_letter = chr(ord('A') + (column_index-1) % 26) + column_letter
        column_index = int((column_index-1) / 26)
    return column_letter

def getExcelData(workbook, sheetName):
    global Separation_Char
    global EmptySpace_Char
    pxl_doc = openpyxl.load_workbook(workbook, data_only = True)
    sheet = pxl_doc[sheetName]
    excelData = {}
    # iterate over columns in the first row
    for i in range(1, sheet.max_column+1):
        # get column letter
        col = get_column_letter(i)
        excelData[sheet.cell(row=1, column=i).value] = sheet.cell(row=2, column=i).value

    ## Used for sinchronizaton over variables within the excel file
    sheet = pxl_doc["CONST"];
    Separation_Char = sheet.cell(row = 1, column = 2).value;
    EmptySpace_Char = sheet.cell(row = 2, column = 2).value;

    return excelData


def deleteImage(image):
    image = image._element;
    image.getparent().remove(image);
    
def deleteShape(shape):
    sp = shape._sp
    sp.getparent().remove(sp)

def getShapeProperties(shape):
    x = shape.left;
    y = shape.top;
    h = shape.height;
    w = shape.width;

    return x, y, h, w;

def fitImage(h_n, w_n, h, w):
    scale = h/h_n;
    h_n = h_n*scale;
    w_n = w_n*scale;

    if w_n > w:
        scale = w/w_n;
        h_n = h_n*scale;
        w_n = w_n*scale;

    return int(h_n), int(w_n);

def fitBackgroundImage(h_n, w_n, h, w):
    scale = w/w_n;
    h_n = h_n*scale;
    w_n = w_n*scale;

    if h_n < h:
        scale = h/h_n;
        h_n = h_n*scale;
        w_n = w_n*scale;
        

    return int(h_n), int(w_n);

def compressImage(f_name, dpi = 96):
    im = Image.open("{}".format(f_name));
    f_name_only = f_name[0:f_name.find(".")];
    im.save("{}".format(f_name), dpi=(dpi,dpi));

    return;

def putImage(slide, shape, source, fmt):
    global global_AUX
    x, y, h, w = getShapeProperties(shape)

    img = source + "." + fmt;

    if os.path.exists(img):
        if fmt == 'jpg':
            compressImage(img);
                
        slide.shapes.add_picture(img, x, y);
            

        if source.find('logo') > 0:
            x, y, h, w = getShapeProperties(slide.shapes[-1]);
            
            logo = source.split('-');
            logo = logo[0];
            logo_file = open("{}-width.txt".format(logo));
            global_AUX = logo_file;
            logo_width = float(logo_file.read());

            scale = logo_width/w.inches;
            
            slide.shapes[-1].width = int(slide.shapes[-1].width*scale);
            slide.shapes[-1].height = int(slide.shapes[-1].height*scale);

            x, y, h, w = getShapeProperties(slide.shapes[-1]);
            
            slide.shapes[-1].top = int(y - h/2);

        elif source.find('background') > 0:
            ## Could compress  background image before processing it
            
            slide.shapes[-1].height, slide.shapes[-1].width = fitBackgroundImage(slide.shapes[-1].height, slide.shapes[-1].width, h, w);
            slide.shapes[-1].top = h - slide.shapes[-1].height;
            
        else:
            slide.shapes[-1].height, slide.shapes[-1].width = fitImage(slide.shapes[-1].height, slide.shapes[-1].width, h, w);

        image = slide.shapes[-1];
        
        slide.shapes._spTree.insert(2, slide.shapes[-1]._element);

    else:
        image = None;
        printWARN("WARNING: Path <{}> does not exists".format(img));

    deleteShape(shape);

    return image;

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def formatOPNTable(slide):
    global global_AUX
    for index, cell in enumerate(iter_cells(slide.shapes[-1].table)):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if index < 3:
                    run.font.size = Pt(16)
                else:
                    run.font.size = Pt(12)
                run.font.name = "Arrow Display"
                
def drawOPNTable(slide, shape, OPN):
    global MAX_ROWS
    global Separation_Char
    x, y, h, w = getShapeProperties(shape)
    OPNs = OPN.split(Separation_Char);

    if len(OPNs) > MAX_ROWS:
        printWARN("WARNING: More than <{0}> P/Ns were identified, taking the first <{0}> elements to preserve the shape".format(MAX_ROWS))
        OPNs = OPNs[0:MAX_ROWS];
    
    slide.shapes.add_table(rows = len(OPNs),
                           cols = 3,
                           left = x,
                           top = y,
                           width = w,
                           height = int(len(OPNs) + 1*h/MAX_ROWS));

def writeOPNTableHeader(slide):
    slide.shapes[-1].table.cell(0, 0).text = "OPN";
    slide.shapes[-1].table.cell(0, 1).text = "Description";
    slide.shapes[-1].table.cell(0, 2).text = "Package";
    

def writeOPNTable(slide, col, data_in):
    global Separation_Char
    data = data_in.split(Separation_Char);

    if len(data) > MAX_ROWS:
        data = data[0:MAX_ROWS];

    for i, val in enumerate(data):
        slide.shapes[-1].table.cell(i, col).text = val;

def writeOPNLinkTable(slide, col, data_in):
    global Separation_Char
    global EmptySpace_Char
    global global_AUX
    
    if not type(data_in) == type(None):
        data = data_in.split(Separation_Char);

        if not len(data) == 0:
            print("PROGRESS: Adding Links to table")
            if len(data) > (MAX_ROWS -1):
                data = data[0:MAX_ROWS-1];

            for i, val in enumerate(data):
                if val == EmptySpace_Char:
                    continue

                aux = slide.shapes[-1].table.cell(i+1, col).text;
                slide.shapes[-1].table.cell(i+1, col).text = "";
                aux_link = slide.shapes[-1].table.cell(i+1, col).text_frame.paragraphs[0].add_run();
                aux_link.text = aux;
                aux_link.hyperlink.address = val;

def verifyLinkLanguage(URL, Header):
    try:
        new_URL = URL.replace("/es-mx/", "/en/");
    except:
        new_URL = URL;
        printWARN("WARNING: Could not modify the Language on the following text <{}> from <{}>".format(URL, Header));

    return new_URL;

def putData(slide, shape, name, data):
    global global_AUX
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            print("PROGRESS: Working with <{}>".format(run.text))
            aux = run.text;

            if not type(data[name]) == type(None):
                if name.find("Image") > 0:
                    if name.find("Logo") >= 0:
                        sourceImage = "images//logo//{}".format(data[name][0:data[name].rfind(".")]);

                    elif name.find("Figure") >= 0:
                        sourceImage = "images//figures//{}".format(data[name][0:data[name].rfind(".")]);

                    elif name.find("Background") >= 0:
                        sourceImage = "images//background//{}".format(data[name][0:data[name].rfind(".")]);

                    elif name.find("App") >= 0:
                        sourceImage = "images//app//{}".format(data[name][0:data[name].rfind(".")]);

                    else:
                        printWARN("WARNING: No folder found for <{}>".format(data[name][0:data[name].rfind(".")]))
                        
                    if data[name].find("svg") > 0:
                        convertSVG2PNG(sourceImage);
                        shape = putImage(slide,
                                         shape,
                                         sourceImage,
                                         "png");

                    elif data[name].find("png") > 0:
                        shape = putImage(slide,
                                         shape,
                                         sourceImage,
                                         "png");

                    elif data[name].find("jpg") > 0:
                        shape = putImage(slide,
                                     shape,
                                     sourceImage,
                                     "jpg");
                        
                    else:
                        printWARN("WARNING: No format for <{}> detected".format(name));

                elif name.find("Table") > 0:

                    if data['OPNTable'].find("Y") == 0:
                        drawOPNTable(slide,
                                     shape,
                                     data['OPNTableColumn1']);

                        writeOPNTable(slide, 0, data['OPNTableColumn1']);
                        writeOPNLinkTable(slide, 0, verifyLinkLanguage(data['OPNTableColumn1Link'], 'OPNTableColumn1Link'));
                        writeOPNTable(slide, 1, data['OPNTableColumn2']);
    ##                    writeOPNLinkTable(slide, 1, verifyLinkLanguage(data['OPNTableColumn2Link'], 'OPNTableColumn2Link'));
                        writeOPNTable(slide, 2, data['OPNTableColumn3']);
                        writeOPNLinkTable(slide, 2, verifyLinkLanguage(data['OPNTableColumn3Link'], 'OPNTableColumn3Link'));
                        formatOPNTable(slide)

                    deleteShape(shape);


                elif name.find("Link") >= 0:
                    aux_mask = aux + "Mask";
                    run.text = "{}".format(data[aux_mask]);
                    run.hyperlink.address = verifyLinkLanguage(data[aux], name);

                elif name.find("Text") >= 0:
                    run.text = data[aux];

                else:
                    max_size = int(run.font.size / 12700);
                    run.text = data[aux];
                    shape.text_frame.fit_text(font_file = "ArrowDisplay_Md.ttf", max_size = max_size);

                    for run in shape.text_frame.paragraphs[0].runs:
                        run.font.name = 'Arrow Display';

            else:
                printWARN("WARNING: Object <{}> has an invalid input <{}>".format(aux, data[name]));
                deleteShape(shape);
                shape = None;                

    if not shape == None:
        x, y, h, w = getShapeProperties(shape);
    else:
        x = None;
        y = None;
        h = None;
        w = None;
        
    return x, y, h, w, shape;






data = getExcelData("./NPI_TEMPLATE_FILL_Test.xlsx", 'TEMPLATE_1_FILL')

prs = Presentation("Template-{}.pptx".format(data['Template']))
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
dict_slide = {};
dict_template = {};
for slide in prs.slides:
    print("PROGRESS: Starting with slide number <{}>".format(prs.slides.index(slide)))
    if not prs.slides.index(slide) in dict_slide:
        dict_template[prs.slides.index(slide)] = {};
        dict_template[prs.slides.index(slide)]['object'] = slide;
        dict_template[prs.slides.index(slide)]['info'] = {};
        
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
##                print(run.text)
                aux = run.text;
                aux = aux.split('-');

                if aux[0] in data:
                    if not aux[0] in dict_template[prs.slides.index(slide)]['info']:
                        print("PROGRESS: Creating index <{}>".format(aux[0]))
                        dict_template[prs.slides.index(slide)]['info'][aux[0]] = {};
                        if len(aux) > 1:
                            dict_template[prs.slides.index(slide)]['info'][aux[0]]['Multiples'] = True;
                            dict_template[prs.slides.index(slide)]['info'][aux[0]]['objects'] = {};
                        else:
                            dict_template[prs.slides.index(slide)]['info'][aux[0]]['Multiples'] = False;
                            dict_template[prs.slides.index(slide)]['info'][aux[0]]['objects'] = {};
                            dict_template[prs.slides.index(slide)]['info'][aux[0]]['objects']['0'] = shape;

                    if dict_template[prs.slides.index(slide)]['info'][aux[0]]['Multiples']:
                        dict_template[prs.slides.index(slide)]['info'][aux[0]]['objects'][aux[-1]] = shape;
                else:
                    printWARN("WARNING: Index <{}> is not inside EXCEL file".format(aux[0]))
                    

for i in dict_template:
    for j in dict_template[i]['info']:
        if dict_template[i]['info'][j]['Multiples']:
            MaxArea = 0;
            list_shapeProp = []
            
            for shape_num in dict_template[i]['info'][j]['objects']:
                list_shapeProp.append([])
                x, y, h, w, dict_template[i]['info'][j]['objects'][shape_num] = putData(dict_template[i]['object'], dict_template[i]['info'][j]['objects'][shape_num], j, data);

                
                if dict_template[i]['info'][j]['objects'][shape_num] == None:
                    printWARN("WARNNIG: <{}> had a None flag, omiting object".format(j))
                    continue;

                list_shapeProp[-1].append(h);
                list_shapeProp[-1].append(w);
                list_shapeProp[-1].append(shape_num);

            if len(list_shapeProp[-1]) > 0:
                for shapeProp in list_shapeProp:
                    area = shapeProp[0]*shapeProp[1];
                    if MaxArea < area:
                        MaxArea = area;
                        i_MaxArea = shapeProp[2];

                for shapeProp in list_shapeProp:
                    if not shapeProp[2] == i_MaxArea:
                        print('Deleting <{}-{}>'.format(j, shapeProp[2]))
                        deleteImage(dict_template[i]['info'][j]['objects'][shapeProp[2]])

            
        else:
            for shape_num in dict_template[i]['info'][j]['objects']:
                x, y, h, w, dict_template[i]['info'][j]['objects'][shape_num] = putData(dict_template[i]['object'], dict_template[i]['info'][j]['objects'][shape_num], j, data);


Prs_Title = "NPI-{}-{}.pptx".format(data['Supplier'], data['PartNumber']);
Prs_Title = Prs_Title.replace("/", "-");

print("Saving <{}>".format(Prs_Title));
prs.save(Prs_Title)
